"""
RAG (Retrieval-Augmented Generation) service for document processing.
Handles PDF and PowerPoint files with chunking, embedding, and semantic search.
"""
import hashlib
import io
import os
import uuid
from datetime import datetime
from typing import Any, Dict, List, Optional

from motor.motor_asyncio import AsyncIOMotorDatabase
from pydantic import BaseModel


class DocumentChunk(BaseModel):
    """A chunk of text from a document."""
    id: str
    document_id: str
    content: str
    page_number: Optional[int] = None
    chunk_index: int
    metadata: Dict[str, Any] = {}


class DocumentOut(BaseModel):
    """Document output schema."""
    id: str
    room_id: str
    filename: str
    file_type: str  # pdf, pptx, docx
    file_size: int
    uploaded_by: str
    uploaded_by_name: Optional[str] = None
    chunk_count: int
    summary: Optional[str] = None
    created_at: str
    
    class Config:
        from_attributes = True


class RAGService:
    """Service for RAG document operations."""
    
    def __init__(self, db: AsyncIOMotorDatabase):
        """Initialize RAG service."""
        self.db = db
        self.documents_collection = db.documents
        self.chunks_collection = db.document_chunks
        self.embeddings_collection = db.embeddings
        
        # Chunk configuration
        self.chunk_size = 1000  # characters
        self.chunk_overlap = 200  # characters overlap between chunks
    
    def _chunk_text(self, text: str, chunk_size: int = None, overlap: int = None) -> List[str]:
        """
        Split text into overlapping chunks.
        
        Args:
            text: Text to chunk
            chunk_size: Size of each chunk
            overlap: Overlap between chunks
            
        Returns:
            List of text chunks
        """
        chunk_size = chunk_size or self.chunk_size
        overlap = overlap or self.chunk_overlap
        
        if len(text) <= chunk_size:
            return [text] if text.strip() else []
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # Try to find a natural break point (sentence end, paragraph)
            if end < len(text):
                # Look for sentence end within last 100 chars
                for sep in ['. ', '.\n', '\n\n', '\n', ' ']:
                    last_sep = text.rfind(sep, start + chunk_size - 100, end)
                    if last_sep > start:
                        end = last_sep + len(sep)
                        break
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            start = end - overlap
        
        return chunks
    
    async def extract_text_from_pdf(self, file_content: bytes) -> List[Dict[str, Any]]:
        """
        Extract text from PDF file.
        
        Args:
            file_content: PDF file bytes
            
        Returns:
            List of dicts with page_number and text
        """
        try:
            from pypdf import PdfReader
            
            pdf_file = io.BytesIO(file_content)
            reader = PdfReader(pdf_file)
            
            pages = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append({
                        "page_number": i + 1,
                        "text": text
                    })
            
            return pages
        except Exception as e:
            print(f"Error extracting PDF text: {e}")
            return []
    
    async def extract_text_from_pptx(self, file_content: bytes) -> List[Dict[str, Any]]:
        """
        Extract text from PowerPoint file.
        
        Args:
            file_content: PPTX file bytes
            
        Returns:
            List of dicts with slide_number and text
        """
        try:
            from pptx import Presentation
            
            pptx_file = io.BytesIO(file_content)
            prs = Presentation(pptx_file)
            
            slides = []
            for i, slide in enumerate(prs.slides):
                texts = []
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
                    
                    # Handle tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text:
                                    texts.append(cell.text)
                
                combined_text = "\n".join(texts)
                if combined_text.strip():
                    slides.append({
                        "page_number": i + 1,
                        "text": combined_text
                    })
            
            return slides
        except Exception as e:
            print(f"Error extracting PPTX text: {e}")
            return []
    
    async def generate_embedding(self, text: str) -> List[float]:
        """
        Generate embedding for text using Gemini.
        
        Args:
            text: Text to embed
            
        Returns:
            Embedding vector
        """
        try:
            from google import genai
            
            client = genai.Client()
            
            # Use text-embedding model
            result = client.models.embed_content(
                model="models/text-embedding-004",
                contents=text[:8000]  # Limit text length
            )
            
            return result.embeddings[0].values
        except Exception as e:
            print(f"Error generating embedding: {e}")
            # Return a zero vector as fallback
            return [0.0] * 768
    
    async def process_document(
        self,
        room_id: str,
        filename: str,
        file_content: bytes,
        file_type: str,
        user_id: str,
        username: str = None
    ) -> Optional[DocumentOut]:
        """
        Process and store a document with RAG capabilities.
        
        Args:
            room_id: Room ID
            filename: Original filename
            file_content: File bytes
            file_type: File type (pdf, pptx)
            user_id: Uploader user ID
            username: Uploader username
            
        Returns:
            DocumentOut or None
        """
        doc_id = str(uuid.uuid4())
        now = datetime.utcnow()
        
        # Extract text based on file type
        pages = []
        if file_type == "pdf":
            pages = await self.extract_text_from_pdf(file_content)
        elif file_type in ["pptx", "ppt"]:
            pages = await self.extract_text_from_pptx(file_content)
        else:
            return None
        
        if not pages:
            return None
        
        # Combine all text for summary
        all_text = "\n\n".join([p["text"] for p in pages])
        
        # Generate summary using Gemini
        summary = await self._generate_summary(all_text, filename)
        
        # Create chunks with embeddings
        chunks = []
        chunk_index = 0
        
        for page in pages:
            page_chunks = self._chunk_text(page["text"])
            
            for chunk_text in page_chunks:
                chunk_id = str(uuid.uuid4())
                
                # Generate embedding for chunk
                embedding = await self.generate_embedding(chunk_text)
                
                chunk = {
                    "id": chunk_id,
                    "document_id": doc_id,
                    "room_id": room_id,
                    "content": chunk_text,
                    "page_number": page.get("page_number"),
                    "chunk_index": chunk_index,
                    "embedding": embedding,
                    "created_at": now
                }
                
                chunks.append(chunk)
                chunk_index += 1
        
        # Store document metadata
        doc = {
            "id": doc_id,
            "room_id": room_id,
            "filename": filename,
            "file_type": file_type,
            "file_size": len(file_content),
            "uploaded_by": user_id,
            "uploaded_by_name": username,
            "chunk_count": len(chunks),
            "summary": summary,
            "created_at": now,
            "updated_at": now
        }
        
        await self.documents_collection.insert_one(doc)
        
        # Store chunks with embeddings
        if chunks:
            await self.chunks_collection.insert_many(chunks)
        
        return DocumentOut(
            id=doc_id,
            room_id=room_id,
            filename=filename,
            file_type=file_type,
            file_size=len(file_content),
            uploaded_by=user_id,
            uploaded_by_name=username,
            chunk_count=len(chunks),
            summary=summary,
            created_at=now.isoformat()
        )
    
    async def _generate_summary(self, text: str, filename: str) -> str:
        """Generate a summary of the document."""
        try:
            from google import genai
            
            client = genai.Client()
            
            prompt = f"""Summarize this document titled "{filename}" in 2-3 sentences. 
            Focus on the main topic and key points.
            
            Document text:
            {text[:4000]}
            """
            
            response = client.models.generate_content(
                model="gemini-2.0-flash-lite",
                contents=prompt
            )
            
            return response.text.strip()
        except Exception as e:
            print(f"Error generating summary: {e}")
            return f"Document: {filename}"
    
    async def get_room_documents(self, room_id: str) -> List[DocumentOut]:
        """Get all documents for a room."""
        cursor = self.documents_collection.find({"room_id": room_id}).sort("created_at", -1)
        docs = await cursor.to_list(length=100)
        
        return [
            DocumentOut(
                id=d["id"],
                room_id=d["room_id"],
                filename=d["filename"],
                file_type=d["file_type"],
                file_size=d["file_size"],
                uploaded_by=d["uploaded_by"],
                uploaded_by_name=d.get("uploaded_by_name"),
                chunk_count=d["chunk_count"],
                summary=d.get("summary"),
                created_at=d["created_at"].isoformat() if isinstance(d["created_at"], datetime) else d["created_at"]
            )
            for d in docs
        ]
    
    async def semantic_search(
        self,
        room_id: str,
        query: str,
        limit: int = 5
    ) -> List[Dict[str, Any]]:
        """
        Search documents in a room using semantic similarity.
        
        Args:
            room_id: Room ID
            query: Search query
            limit: Max results
            
        Returns:
            List of matching chunks with scores
        """
        # Generate query embedding
        query_embedding = await self.generate_embedding(query)
        
        # Get all chunks for the room
        chunks = await self.chunks_collection.find(
            {"room_id": room_id}
        ).to_list(length=1000)
        
        if not chunks:
            return []
        
        # Calculate cosine similarity
        results = []
        for chunk in chunks:
            chunk_embedding = chunk.get("embedding", [])
            if chunk_embedding:
                similarity = self._cosine_similarity(query_embedding, chunk_embedding)
                results.append({
                    "chunk_id": chunk["id"],
                    "document_id": chunk["document_id"],
                    "content": chunk["content"],
                    "page_number": chunk.get("page_number"),
                    "similarity": similarity
                })
        
        # Sort by similarity and return top results
        results.sort(key=lambda x: x["similarity"], reverse=True)
        return results[:limit]
    
    def _cosine_similarity(self, a: List[float], b: List[float]) -> float:
        """Calculate cosine similarity between two vectors."""
        import math
        
        if len(a) != len(b):
            return 0.0
        
        dot_product = sum(x * y for x, y in zip(a, b))
        magnitude_a = math.sqrt(sum(x * x for x in a))
        magnitude_b = math.sqrt(sum(x * x for x in b))
        
        if magnitude_a == 0 or magnitude_b == 0:
            return 0.0
        
        return dot_product / (magnitude_a * magnitude_b)
    
    async def ask_document(
        self,
        room_id: str,
        question: str,
        document_id: str = None
    ) -> str:
        """
        Answer a question using RAG on documents.
        
        Args:
            room_id: Room ID
            question: User question
            document_id: Optional specific document to query
            
        Returns:
            AI-generated answer
        """
        # Get relevant chunks
        relevant_chunks = await self.semantic_search(room_id, question, limit=5)
        
        if not relevant_chunks:
            return "I couldn't find relevant information in the uploaded documents."
        
        # Build context from chunks
        context_parts = []
        for i, chunk in enumerate(relevant_chunks):
            page_info = f" (Page {chunk['page_number']})" if chunk.get('page_number') else ""
            context_parts.append(f"[Source {i+1}{page_info}]: {chunk['content']}")
        
        context = "\n\n".join(context_parts)
        
        # Generate answer using Gemini
        try:
            from google import genai
            
            client = genai.Client()
            
            prompt = f"""Based on the following document excerpts, answer the question.
            If the answer cannot be found in the excerpts, say so.
            
            Document excerpts:
            {context}
            
            Question: {question}
            
            Answer:"""
            
            response = client.models.generate_content(
                model="gemini-2.0-flash-lite",
                contents=prompt
            )
            
            return response.text.strip()
        except Exception as e:
            print(f"Error generating answer: {e}")
            return "I encountered an error while processing your question."
    
    async def delete_document(self, document_id: str) -> bool:
        """Delete a document and its chunks."""
        # Delete chunks
        await self.chunks_collection.delete_many({"document_id": document_id})
        
        # Delete document
        result = await self.documents_collection.delete_one({"id": document_id})
        
        return result.deleted_count > 0
